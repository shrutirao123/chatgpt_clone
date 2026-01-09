import json
import os
import time
import io
import base64
import re  # <--- NEW: Added for username validation
import requests
import pandas as pd
import concurrent.futures
from threading import Thread
from PIL import Image, ImageOps

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.models import User
from django.contrib.auth.decorators import login_required
from django.views.decorators.csrf import csrf_exempt
from django.http import JsonResponse, StreamingHttpResponse
from django.core.cache import cache
from django.core.files.storage import default_storage
from django.conf import settings
from django.db import close_old_connections

from rest_framework.decorators import api_view
from rest_framework.response import Response

# Document Loaders
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

# Local Imports
from .models import Conversation, ChatHistory
from .rag_utils import add_document, retrieve_relevant_chunks, delete_index

# --- Config ---
AI_SERVER_URL = "http://localhost:11434/api/generate"
AI_MODEL = "mistral-nemo:12b"
VISION_MODEL = "llama3.2-vision:11b"

# Performance & Constraints
MAX_IMAGE_SIZE = (512, 512)
JPEG_QUALITY = 65
RAG_TOP_K = 3
SAVE_INTERVAL_WORDS = 25  # Save DB less frequently to reduce I/O
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
MAX_TEXT_LENGTH = 30000

# Limits for file parsers
MAX_PDF_PAGES = 30
MAX_DOC_PARAGRAPHS = 50
MAX_PPT_SLIDES = 20
MAX_CSV_ROWS = 500


# ==========================================
#              AUTHENTICATION
# ==========================================

def login_view(request):
    if request.user.is_authenticated:
        return redirect('chat_ui_base')
        
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')
        user = authenticate(request, username=username, password=password)
        if user:
            login(request, user)
            return redirect('chat_ui_base')
        return render(request, 'chat/login.html', {'error': 'Invalid credentials'})
    return render(request, 'chat/login.html')


def register_view(request):
    if request.user.is_authenticated:
        return redirect('chat_ui_base')

    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')
        email = request.POST.get('email') # Email bhi fetch kar rahe hain

        # --- VALIDATION START ---

        # 1. Username Space Check
        if ' ' in username:
            return render(request, 'chat/register.html', {'error': "Username cannot contain spaces. Use underscores."})

        # 2. Username Strict Characters (Only a-z, 0-9, _)
        if not re.match(r'^[a-zA-Z0-9_]+$', username):
            return render(request, 'chat/register.html', {'error': "Invalid username. Use only letters, numbers, and underscores."})

        # 3. Password Special Character Check (NEW)
        # Check agar kam se kam 1 special character hai
        if not re.search(r'[!@#$%^&*(),.?":{}|<>]', password):
            return render(request, 'chat/register.html', {'error': "Password must contain at least one special character (!@#$%)."})

        # 4. Password Length Check (Optional but good)
        if len(password) < 6:
            return render(request, 'chat/register.html', {'error': "Password must be at least 6 characters long."})

        # --- VALIDATION END ---

        if User.objects.filter(username=username).exists():
            return render(request, 'chat/register.html', {'error': 'Username already exists'}) 

        try:
            # Email field bhi save kar rahe hain agar model mein hai
            user = User.objects.create_user(username=username, password=password, email=email)
            login(request, user)
            return redirect('chat_ui_base')
        except Exception as e:
            return render(request, 'chat/register.html', {'error': f'Registration failed: {e}'})

    return render(request, 'chat/register.html')

def logout_view(request):
    logout(request)
    return redirect('login')


# ==========================================
#                 CHAT UI
# ==========================================

@login_required
def chat_ui(request, conversation_id=None):
    conversations = Conversation.objects.filter(user=request.user, is_archived=False).order_by('-created_at')[:30]
    archived_conversations = Conversation.objects.filter(user=request.user, is_archived=True).order_by('-created_at')[:20]

    if not conversation_id:
        # Check if the most recent conversation is empty (New Chat), reuse it to prevent spamming DB
        last_conv = conversations.first()
        if last_conv and last_conv.title == "New Chat" and not last_conv.messages.exists():
            return redirect('chat_ui_conversation', conversation_id=last_conv.id)
            
        conversation = Conversation.objects.create(user=request.user, title="New Chat")
        return redirect('chat_ui_conversation', conversation_id=conversation.id)

    conversation = get_object_or_404(Conversation, id=conversation_id, user=request.user)
    messages = conversation.messages.all().order_by('created_at')

    return render(request, 'chat/chat.html', {
        "conversations": conversations,
        "archived_conversations": archived_conversations,
        "conversation": conversation,
        "messages": messages
    })


# ==========================================
#              API ENDPOINTS
# ==========================================

@csrf_exempt
@login_required
def rename_conversation(request, conv_id):
    if request.method != "POST":
        return JsonResponse({"error": "Method not allowed"}, status=405)
        
    conversation = get_object_or_404(Conversation, id=conv_id, user=request.user)
    try:
        data = json.loads(request.body)
        new_title = data.get("title", "").strip()[:100]
        if not new_title:
            return JsonResponse({"error": "Title cannot be empty"}, status=400)
        
        conversation.title = new_title
        conversation.save(update_fields=['title'])
        return JsonResponse({"success": True, "title": conversation.title})
    except json.JSONDecodeError:
        return JsonResponse({"error": "Invalid JSON"}, status=400)


@csrf_exempt
@login_required
def toggle_archive_conversation(request, conv_id):
    if request.method != "POST":
        return JsonResponse({"error": "Method not allowed"}, status=405)

    conversation = get_object_or_404(Conversation, id=conv_id, user=request.user)
    conversation.is_archived = not conversation.is_archived
    conversation.save(update_fields=['is_archived'])
    
    return JsonResponse({
        "success": True, 
        "is_archived": conversation.is_archived,
        "message": "Archived" if conversation.is_archived else "Unarchived"
    })


@csrf_exempt
@login_required
def delete_conversation(request, conv_id):
    if request.method != "DELETE":
        return JsonResponse({"error": "Method not allowed"}, status=405)

    conversation = get_object_or_404(Conversation, id=conv_id, user=request.user)
    
    def delete_async():
        # Ensure DB connection is closed in thread
        try:
            delete_index(conv_id)
        except Exception as e:
            print(f"Error deleting index: {e}")
        finally:
            close_old_connections()
    
    Thread(target=delete_async, daemon=True).start()
    conversation.delete()
    return JsonResponse({"success": True, "message": "Conversation deleted."})


@login_required
@api_view(['POST'])
def new_conversation(request):
    conversation = Conversation.objects.create(user=request.user, title="New Chat")
    return Response({"conversation_id": conversation.id}, status=201)


# ==========================================
#           HELPER FUNCTIONS
# ==========================================

def is_image_file(filename):
    return filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'))


def get_file_icon(filename):
    ext_map = {
        '.pdf': '📕',
        '.docx': '📘', '.doc': '📘',
        '.pptx': '📊', '.ppt': '📊',
        '.xlsx': '📈', '.xls': '📈',
        '.csv': '📋',
        '.txt': '📄', '.md': '📄'
    }
    for ext, icon in ext_map.items():
        if filename.lower().endswith(ext):
            return icon
    return '📎'


def format_file_size(bytes_size):
    for unit in ['B', 'KB', 'MB', 'GB']:
        if bytes_size < 1024.0:
            return f"{bytes_size:.1f} {unit}"
        bytes_size /= 1024.0
    return f"{bytes_size:.1f} TB"


def process_image(uploaded_file):
    """Resizes and encodes image to Base64, handling EXIF rotation."""
    try:
        image = Image.open(uploaded_file)
        # Fix orientation (e.g. from phones)
        image = ImageOps.exif_transpose(image)
        
        image.thumbnail(MAX_IMAGE_SIZE, Image.Resampling.LANCZOS)
        
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='JPEG', quality=JPEG_QUALITY, optimize=True)
        return base64.b64encode(img_byte_arr.getvalue()).decode('utf-8')
    except Exception as e:
        print(f"Error processing image: {e}")
        return None


def save_uploaded_image(uploaded_file, conversation_id):
    try:
        upload_dir = os.path.join('chat_images', str(conversation_id))
        # Ensure unique filename
        filename = f"{int(time.time())}_{uploaded_file.name.replace(' ', '_')}"
        file_path = os.path.join(upload_dir, filename)
        full_path = default_storage.save(file_path, uploaded_file)
        return settings.MEDIA_URL + full_path
    except Exception as e:
        print(f"Error saving image: {e}")
        return None


def get_conversation_images(conversation_id):
    return cache.get(f"conv_images_{conversation_id}", [])


def store_conversation_image(conversation_id, base64_image, image_url):
    cache_key = f"conv_images_{conversation_id}"
    images = cache.get(cache_key, [])
    images.append({"base64": base64_image, "url": image_url})
    # Keep last 5 images context
    if len(images) > 5:
        images = images[-5:]
    cache.set(cache_key, images, timeout=86400)


def extract_text_from_file(uploaded_file):
    """Extracts text from various file formats."""
    name = uploaded_file.name.lower()
    text = ""
    
    if uploaded_file.size > MAX_FILE_SIZE:
        return f"[File too large. Max size: {format_file_size(MAX_FILE_SIZE)}]"

    try:
        if name.endswith(".pdf"):
            reader = PdfReader(uploaded_file)
            pages_to_read = min(len(reader.pages), MAX_PDF_PAGES)
            text = "\n".join([page.extract_text() or "" for page in reader.pages[:pages_to_read]])
            if len(reader.pages) > MAX_PDF_PAGES:
                text += f"\n\n[Truncated: Read {MAX_PDF_PAGES}/{len(reader.pages)} pages]"

        elif name.endswith(('.docx', '.doc')):
            doc = DocxDocument(uploaded_file)
            paragraphs = doc.paragraphs[:MAX_DOC_PARAGRAPHS]
            text = "\n".join([para.text for para in paragraphs if para.text.strip()])
            if len(doc.paragraphs) > MAX_DOC_PARAGRAPHS:
                text += f"\n\n[Truncated: Read {MAX_DOC_PARAGRAPHS}/{len(doc.paragraphs)} paragraphs]"

        elif name.endswith(('.ppt', '.pptx')):
            prs = Presentation(uploaded_file)
            slides_to_read = min(len(prs.slides), MAX_PPT_SLIDES)
            slides_text = []
            for i, slide in enumerate(list(prs.slides)[:slides_to_read]):
                shape_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if shape_text:
                    slides_text.append(f"\n--- Slide {i+1} ---\n" + "\n".join(shape_text))
            text = "".join(slides_text)
            if len(prs.slides) > MAX_PPT_SLIDES:
                text += f"\n\n[Truncated: Read {MAX_PPT_SLIDES}/{len(prs.slides)} slides]"

        elif name.endswith(".csv"):
            df = pd.read_csv(uploaded_file, dtype=str, keep_default_na=False, nrows=MAX_CSV_ROWS)
            text = df.to_string(index=False, max_rows=MAX_CSV_ROWS)

        elif name.endswith(('.xls', '.xlsx')):
            xls = pd.read_excel(uploaded_file, sheet_name=None, dtype=str, nrows=MAX_CSV_ROWS)
            for sheet_name, df in list(xls.items())[:3]:
                text += f"\n\nSheet: {sheet_name}\n{df.fillna('').to_string(index=False, max_rows=MAX_CSV_ROWS)}"

        else:
            # Fallback for text files
            content = uploaded_file.read(50 * 1024) # Read first 50KB
            text = content.decode(errors="ignore")

    except Exception as e:
        return f"[Error reading file: {str(e)[:150]}]"

    if len(text) > MAX_TEXT_LENGTH:
        text = text[:MAX_TEXT_LENGTH] + "\n\n[Content truncated for speed]"

    return text.strip()


def index_document_async(conversation_id, text, filename, file_size):
    """Indexes document in FAISS/Vector DB in a background thread."""
    def do_index():
        try:
            text_to_index = text[:MAX_TEXT_LENGTH]
            add_document(conversation_id, text_to_index)
            print(f"✓ Indexed: {filename}")
        except Exception as e:
            print(f"✗ Indexing error for {filename}: {e}")
        finally:
            close_old_connections()
    
    Thread(target=do_index, daemon=True).start()


def construct_prompt(conversation_id, prompt, uploaded_file=None):
    """Builds the final system/user prompt including RAG context."""
    full_prompt = ""
    
    # RAG Context
    if not uploaded_file:
        context_chunks = retrieve_relevant_chunks(conversation_id, prompt, top_k=RAG_TOP_K)
        if context_chunks:
            full_prompt += "Stored Knowledge (Use this if relevant):\n" + "\n\n".join(context_chunks) + "\n\n"
            
    full_prompt += f"Question: {prompt}\n\n"
    
    full_prompt += """Instructions:
- Format response in clean Markdown.
- Use **bold** for key terms.
- Avoid unnecessary conversational filler.
"""
    return full_prompt


# ==========================================
#           CORE AI LOGIC
# ==========================================

def handle_uploaded_file_logic(conversation, uploaded_file, prompt):
    """
    Handles saving file, extracting text, creating user messages in DB.
    Returns: (text_content, base64_image_if_any, is_image_bool)
    """
    is_image = is_image_file(uploaded_file.name)
    conversation_id = conversation.id
    base64_image = None
    extracted_text = None

    if is_image:
        with concurrent.futures.ThreadPoolExecutor(max_workers=2) as executor:
            future_save = executor.submit(save_uploaded_image, uploaded_file, conversation_id)
            uploaded_file.seek(0)
            future_process = executor.submit(process_image, uploaded_file)
            
            image_url = future_save.result()
            base64_image = future_process.result()

        if base64_image:
            store_conversation_image(conversation_id, base64_image, image_url)

        if image_url:
            user_content = f'<img src="{image_url}" alt="uploaded image" style="max-width:300px; border-radius:8px;"><br>{prompt if prompt else "Describe this image."}'
            ChatHistory.objects.create(conversation=conversation, sender="user", content=user_content)

    else:
        # Document handling
        file_icon = get_file_icon(uploaded_file.name)
        file_size = format_file_size(uploaded_file.size)
        
        user_file_content = f'<div style="display: flex; align-items: center; padding: 10px; background: rgba(100,100,100,0.1); border-radius: 8px; margin-bottom: 8px;"><span style="font-size: 24px; margin-right: 10px;">{file_icon}</span><div><strong>{uploaded_file.name}</strong><br><small style="opacity: 0.7;">{file_size}</small></div></div>{prompt if prompt else "Process this file"}'
        ChatHistory.objects.create(conversation=conversation, sender="user", content=user_file_content)

        extracted_text = extract_text_from_file(uploaded_file)
        
        if extracted_text and not extracted_text.startswith("[Error"):
            index_document_async(conversation_id, extracted_text, uploaded_file.name, file_size)
            sys_msg = f"📄 Indexing: {uploaded_file.name} ({file_size})"
            ChatHistory.objects.create(conversation=conversation, sender="system", content=sys_msg)

    return extracted_text, base64_image, is_image


@csrf_exempt
@login_required
def stream_chat_with_ai(request):
    """
    Optimized Streaming Endpoint.
    """
    if request.method != "POST":
        return JsonResponse({"error": "Only POST allowed"}, status=405)

    # 1. Handle Stop Request
    if request.POST.get("stop") == "true":
        conversation_id = request.POST.get("conversation_id")
        if conversation_id:
            cache.set(f"stop_stream_{conversation_id}", True, timeout=10)
        return JsonResponse({"status": "stopping"})

    # 2. Parse Request
    prompt = request.POST.get("message", "").strip()
    conversation_id = request.POST.get("conversation_id")
    uploaded_file = request.FILES.get("file")

    if not conversation_id:
        return JsonResponse({"error": "Missing conversation_id"}, status=400)

    conversation = get_object_or_404(Conversation, id=conversation_id, user=request.user)
    cache.delete(f"stop_stream_{conversation_id}")

    # 3. Handle Files & User Message Creation
    full_prompt = prompt
    images_to_send = None
    ai_model = AI_MODEL

    if uploaded_file:
        text, b64_img, is_img = handle_uploaded_file_logic(conversation, uploaded_file, prompt)
        
        if is_img and b64_img:
            ai_model = VISION_MODEL
            images_to_send = [b64_img]
            full_prompt = prompt or "Describe this image."
        elif text:
            # It's a document
            if not prompt:
                summary_text = text[:8000]
                full_prompt = f"Summarize this document:\n\n{summary_text}"
            else:
                context_text = text[:12000]
                full_prompt = f"Document Context:\n{context_text}\n\nQuestion: {prompt}"
        else:
            return JsonResponse({"error": "Failed to process file"}, status=400)
    else:
        # Text only message
        ChatHistory.objects.create(conversation=conversation, sender="user", content=prompt)
        
        # Check for previous images in context
        stored_images = get_conversation_images(conversation_id)
        if stored_images:
            ai_model = VISION_MODEL
            images_to_send = [img["base64"] for img in stored_images[-1:]] # Just last image for speed
            full_prompt = f"Referencing previous image. {construct_prompt(conversation_id, prompt)}"
        else:
            full_prompt = construct_prompt(conversation_id, prompt)

    # 4. Update Conversation Title if new
    if conversation.title == "New Chat":
        conversation.title = prompt.split('\n')[0][:50] if prompt else "File Upload"
        conversation.save(update_fields=['title'])

    # 5. Create Placeholder Bot Message
    ai_message_history = ChatHistory.objects.create(conversation=conversation, sender="bot", content="")

    # 6. Stream Generator
    def event_stream():
        ai_response_accum = ""
        save_word_count = 0
        buffer = ""
        in_fenced_code = False
        
        payload = {
            "model": ai_model,
            "prompt": full_prompt,
            "stream": True,
            "options": {"num_ctx": 4096, "temperature": 0.7}
        }
        
        if images_to_send:
            payload["images"] = images_to_send

        try:
            with requests.post(AI_SERVER_URL, json=payload, stream=True, timeout=60) as r:
                r.raise_for_status()

                for line in r.iter_lines(decode_unicode=False):
                    # Check stop signal
                    if cache.get(f"stop_stream_{conversation_id}"):
                        cache.delete(f"stop_stream_{conversation_id}")
                        yield "data: [STOPPED]\n\n"
                        break

                    if not line: continue

                    try:
                        data = json.loads(line.decode("utf-8", errors="ignore"))
                        token = data.get("response", "")
                        if not token: continue
                    except ValueError:
                        continue

                    # Buffer logic to clean up tokens
                    if "```" in token:
                        in_fenced_code = not in_fenced_code
                    
                    buffer += token
                    
                    # Flush buffer intelligently (on spaces/punctuation or size)
                    should_flush = len(buffer) > 5 or (not in_fenced_code and buffer[-1] in " \n.,!?:")
                    
                    if should_flush:
                        ai_response_accum += buffer
                        yield f"data: {buffer}\n\n"
                        buffer = ""
                        
                        # Intermediate DB save
                        current_words = len(ai_response_accum.split())
                        if current_words - save_word_count >= SAVE_INTERVAL_WORDS:
                            ChatHistory.objects.filter(pk=ai_message_history.pk).update(content=ai_response_accum)
                            save_word_count = current_words

                # Flush remaining buffer
                if buffer:
                    ai_response_accum += buffer
                    yield f"data: {buffer}\n\n"

                # Final Save
                ChatHistory.objects.filter(pk=ai_message_history.pk).update(content=ai_response_accum)

        except requests.exceptions.ConnectionError:
            error_msg = "Error: AI Server is offline."
            yield f"data: {error_msg}\n\n"
            ChatHistory.objects.filter(pk=ai_message_history.pk).update(content=error_msg)
        except Exception as e:
            error_msg = f"Error: {str(e)[:100]}"
            yield f"data: {error_msg}\n\n"
            ChatHistory.objects.filter(pk=ai_message_history.pk).update(content=error_msg)

    response = StreamingHttpResponse(event_stream(), content_type="text/event-stream")
    response['Cache-Control'] = 'no-cache'
    return response


@csrf_exempt
@login_required
def chat_with_ai(request):
    """
    Standard (Non-Streaming) Endpoint.
    """
    if request.method != 'POST':
        return JsonResponse({"error": "Only POST allowed"}, status=405)

    prompt = request.POST.get("message", "").strip()
    conversation_id = request.POST.get("conversation_id")
    uploaded_file = request.FILES.get("file")

    if not conversation_id:
        return JsonResponse({"error": "conversation_id is required"}, status=400)

    conversation = get_object_or_404(Conversation, id=conversation_id, user=request.user)
    
    # Reuse logic from stream view logic conceptually
    full_prompt = prompt
    images_to_send = None
    ai_model = AI_MODEL

    if uploaded_file:
        text, b64_img, is_img = handle_uploaded_file_logic(conversation, uploaded_file, prompt)
        if is_img and b64_img:
            ai_model = VISION_MODEL
            images_to_send = [b64_img]
        elif text:
            full_prompt = f"Document Context:\n{text[:10000]}\n\nQuestion: {prompt}"
    else:
        ChatHistory.objects.create(conversation=conversation, sender="user", content=prompt)
        full_prompt = construct_prompt(conversation_id, prompt)
    
    # Prepare Payload
    payload = {
        "model": ai_model,
        "prompt": full_prompt,
        "stream": False,
        "options": {"num_ctx": 4096}
    }
    if images_to_send:
        payload["images"] = images_to_send

    try:
        response = requests.post(AI_SERVER_URL, json=payload, timeout=120)
        response.raise_for_status()
        ai_text = response.json().get("response", "No response.")
    except Exception as e:
        ai_text = f"AI Error: {str(e)}"

    ChatHistory.objects.create(conversation=conversation, sender="bot", content=ai_text)
    return JsonResponse({"response": ai_text})